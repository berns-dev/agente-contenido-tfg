"""Extraccion de texto desde PDF/PPTX."""

from __future__ import annotations

import re
from pathlib import Path

from cleaner import clean_extracted_text


_CONSONANT_RUN_RE = re.compile(r"(?i)[bcdfghjklmnñpqrstvwxyz]{5,}")
_KNOWN_WORDS = {
    "de",
    "del",
    "la",
    "el",
    "tema",
    "chapter",
    "section",
    "slide",
    "content",
    "engineering",
    "machine",
    "elements",
}


def _is_mirrored_text(line: str) -> bool:
    """
    Detecta texto extraído en espejo (rotado 180°).
    Heurística: una línea es texto en espejo si al invertirla
    forma palabras reales en español o inglés más que el original.
    Señales simples: secuencias de consonantes sin vocales >4 chars,
    o palabras conocidas al revertir.
    """
    stripped = line.strip()
    if len(stripped) < 4:
        return False
    reversed_line = stripped[::-1]
    vowels = set("aeiouáéíóúAEIOUÁÉÍÓÚ")
    original_vowels = sum(1 for c in stripped if c in vowels)
    reversed_vowels = sum(1 for c in reversed_line if c in vowels)
    if len(stripped) > 0:
        ratio = original_vowels / len(stripped)
        if ratio < 0.1 and reversed_vowels > original_vowels:
            return True

    original_tokens = [t.lower() for t in stripped.split()]
    reversed_tokens = [t.lower() for t in reversed_line.split()]
    original_known = sum(1 for t in original_tokens if t in _KNOWN_WORDS)
    reversed_known = sum(1 for t in reversed_tokens if t in _KNOWN_WORDS)
    if reversed_known > original_known and reversed_known > 0:
        return True

    if _CONSONANT_RUN_RE.search(stripped):
        rev_has_consonant_run = _CONSONANT_RUN_RE.search(reversed_line) is not None
        if not rev_has_consonant_run:
            return True
    return False


def _extract_pdf(path: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            try:
                texto_pagina = page.extract_text() or ""
                print(f"[EXTRACTOR] Página {idx}: {len(texto_pagina)} chars extraídos")
                print(f"[EXTRACTOR] Primeros 200 chars: {repr(texto_pagina[:200])}")
                filtered_lines = [
                    ln for ln in texto_pagina.split("\n") if not _is_mirrored_text(ln)
                ]
                text = clean_extracted_text("\n".join(filtered_lines), path.name).strip()
                if text:
                    parts.append(f"[PAGINA {idx}]\n{text}")
                else:
                    parts.append(f"[PAGINA {idx}]\n[TEXTO_ILEGIBLE]")
            except Exception:
                parts.append(f"[PAGINA {idx}]\n[TEXTO_ILEGIBLE]")
    return "\n\n".join(parts).strip()


def _escape_markdown_table_cell(text: str) -> str:
    """
    Escapa contenido de celda para tablas Markdown con pipes.
    Sin esto, un '|' literal en la celda rompe la tabla y confunde al pipeline.
    """
    return text.replace("\\", "\\\\").replace("|", "\\|")


def _pptx_table_to_markdown(table: object) -> str:
    """
    Convierte una tabla de python-pptx a Markdown tipo GitHub (pipes).
    La primera fila se trata como cabecera: es la convención más habitual en PPTX
    y coincide con el enfoque de MarkItDown (primera fila = th).
    """
    rows = list(getattr(table, "rows", []) or [])
    if not rows:
        return ""
    lines: list[str] = []
    for row_idx, row in enumerate(rows):
        cells = [
            _escape_markdown_table_cell(
                str(getattr(cell, "text", "") or "").strip().replace("\n", " ")
            )
            for cell in row.cells
        ]
        lines.append("| " + " | ".join(cells) + " |")
        if row_idx == 0:
            lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n".join(lines)


def _pptx_collect_shape_blocks(shape: object, title_shape: object | None) -> list[str]:
    """
    Recorre shapes en el orden que expone python-pptx (orden Z de la diapositiva).
    Los grupos se expanden: el texto en grupos anidados no aparece en el iterador
    plano de slide.shapes y se perdería sin recursión.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

    blocks: list[str] = []

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for child in getattr(shape, "shapes", []):
            blocks.extend(_pptx_collect_shape_blocks(child, title_shape))
        return blocks

    # El título de slide ya va en un bloque `#` aparte: no duplicar el mismo shape.
    if title_shape is not None and shape is title_shape:
        return blocks

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.TABLE and getattr(
        shape, "table", None
    ):
        md_table = _pptx_table_to_markdown(shape.table)
        if md_table.strip():
            blocks.append(md_table.strip())
        return blocks

    if getattr(shape, "has_text_frame", False):
        raw = str(getattr(shape, "text", "") or "").strip()
        if not raw:
            return blocks
        # Subtítulo de layout típico: segunda jerarquía sin inferir roles complejos
        # (solo placeholders explícitos de PPTX, evitamos heurísticas frágiles).
        prefix = ""
        try:
            if getattr(shape, "is_placeholder", False):
                ph_type = shape.placeholder_format.type
                subtitle_types = {PP_PLACEHOLDER.SUBTITLE}
                if hasattr(PP_PLACEHOLDER, "VERTICAL_SUBTITLE"):
                    subtitle_types.add(PP_PLACEHOLDER.VERTICAL_SUBTITLE)
                if ph_type in subtitle_types:
                    prefix = "## "
        except (ValueError, AttributeError):
            pass
        blocks.append(prefix + raw)
        return blocks

    return blocks


def _pptx_slide_to_text(slide: object) -> str:
    """
    Ensambla el contenido de una diapositiva: título como H1, cuerpo/tablas en orden,
    notas al final en un bloque propio (trazabilidad docente / guion del profesor).
    """
    blocks: list[str] = []

    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and getattr(title_shape, "text", None):
        title_text = str(title_shape.text).strip()
        if title_text:
            # Un solo H1 por slide: refleja el modelo mental de "título de diapositiva"
            # y alinea el Markdown con la jerarquía visual habitual.
            blocks.append(f"# {title_text}")

    for shape in slide.shapes:
        blocks.extend(_pptx_collect_shape_blocks(shape, title_shape))

    if getattr(slide, "has_notes_slide", False):
        notes_tf = slide.notes_slide.notes_text_frame
        if notes_tf is not None:
            notes_text = str(notes_tf.text or "").strip()
            if notes_text:
                # Bloque separado y estable para chunking/LLM sin mezclarlo con el bullet principal.
                blocks.append("### Notas del presentador\n\n" + notes_text)

    # Doble salto entre bloques: separa tablas Markdown del texto plano sin fusionar filas.
    return "\n\n".join(b for b in blocks if b).strip()


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        try:
            slide_raw_text = _pptx_slide_to_text(slide)
            if slide_raw_text:
                print(f"[EXTRACTOR] Slide {idx}: {len(slide_raw_text)} chars extraídos")
                print(f"[EXTRACTOR] Primeros 200 chars: {repr(slide_raw_text[:200])}")
                filtered_lines = [
                    ln for ln in slide_raw_text.split("\n") if not _is_mirrored_text(ln)
                ]
                slide_text = clean_extracted_text("\n".join(filtered_lines), path.name).strip()
                if slide_text:
                    parts.append(f"[SLIDE {idx}]\n{slide_text}")
                else:
                    parts.append(f"[SLIDE {idx}]\n[TEXTO_ILEGIBLE]")
            else:
                parts.append(f"[SLIDE {idx}]\n[TEXTO_ILEGIBLE]")
        except Exception:
            parts.append(f"[SLIDE {idx}]\n[TEXTO_ILEGIBLE]")
    return "\n\n".join(parts).strip()


def extract_text(file_path: str) -> str:
    """Extrae texto de un PDF o PPTX."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"No existe el archivo: {file_path}")

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".pptx":
        return _extract_pptx(path)

    raise ValueError("Formato no soportado. Usa PDF o PPTX.")
